from magika import Magika
from magika.types import MagikaResult 
from chardet.universaldetector import UniversalDetector
from pathlib import Path

class FileUtil:

    @classmethod    
    def sanitize_text(cls, text: str) -> str:
        # テキストをサニタイズする
        # textが空の場合は空の文字列を返す
        if not text or len(text) == 0:
            return ""
        import re
        # 1. 複数の改行を1つの改行に変換
        text = re.sub(r'\n+', '\n', text)
        # 2. 複数のスペースを1つのスペースに変換
        text = re.sub(r' +', ' ', text)

        return text

    @classmethod
    def identify_type(cls, filename):
        m = Magika()
        # ファイルの種類を判定
        path = Path(filename)
        try:
            res: MagikaResult = m.identify_path(path) # type: ignore
            encoding = None
            if res.dl.is_text:
                encoding = cls.get_encoding(filename)
        except Exception as e:
            print(e)
            return None, None

        return res, encoding

    @classmethod
    def get_encoding(cls, filename):
        # ファイルのbyte列を取得
        # アクセスできない場合は例外をキャッチ
        try:
            with open(filename, "rb") as f:
                # 1KB読み込む
                byte_data = f.read(8192)
                if not byte_data:
                    return None, None
        except Exception as e:
            print(e)
            return None, None
        # エンコーディング判定
        detector = UniversalDetector()
        detector.feed(byte_data)
        detector.close()
        encoding = detector.result['encoding']  
        return encoding

    @classmethod
    def get_mime_type(cls, filename):
        res, encoding = cls.identify_type(filename)
        if res is None:
            return None
        return res.output.mime_type

    @classmethod
    def extract_text_from_file(cls, filename):
        res, encoding = cls.identify_type(filename)
        
        if res is None:
            return None
        print(res.output.mime_type)
        result = None        
        if res.output.mime_type.startswith("text/"):
            result = cls.process_text(filename, res, encoding)

        # application/pdf
        elif res.output.mime_type == "application/pdf":
            result = cls.process_pdf(filename)
            
        # application/vnd.openxmlformats-officedocument.spreadsheetml.sheet
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            result = ExcelUtil.extract_text_from_sheet(filename)
            
        # application/vnd.openxmlformats-officedocument.wordprocessingml.document
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            result = cls.process_docx(filename)
            
        # application/vnd.openxmlformats-officedocument.presentationml.presentation
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            result = cls.process_pptx(filename)
        else:
            print("Unsupported file type: " + res.output.mime_type)

        return cls.sanitize_text(result)

    # application/pdfのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_pdf(cls, filename):
        from pdfminer.high_level import extract_text
        text = extract_text(filename)
        return text

    # application/vnd.openxmlformats-officedocument.wordprocessingml.documentのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_docx(cls, filename):
        import docx
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        doc = docx.Document(filename)
        for para in doc.paragraphs:
            output.write(para.text)
            output.write("\n")
            
        return output.getvalue()

    # application/vnd.openxmlformats-officedocument.presentationml.presentationのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_pptx(cls, filename):
        import pptx
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        prs = pptx.Presentation(filename)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output.write(shape.text)
                    output.write("\n")
        
        return output.getvalue()

    # text/*のファイルを読み込んで文字列として返す関数
    @classmethod
    def process_text(cls, filename, res, encoding):
        result = ""
        if res.output.mime_type == "text/html":
            # text/htmlの場合
            from bs4 import BeautifulSoup
            # テキストを取得
            with open(filename, "rb") as f:
                text_data = f.read()
                soup = BeautifulSoup(text_data, "html.parser")
            result = soup.get_text()

        elif res.output.mime_type == "text/xml":
            # text/xmlの場合
            from bs4 import BeautifulSoup
            # テキストを取得
            with open(filename, "rb") as f:
                text_data = f.read()
                soup = BeautifulSoup(text_data, features="xml")
            result = soup.get_text()

        elif res.output.mime_type == "text/markdown":
            # markdownの場合
            from bs4 import BeautifulSoup
            from markdown import markdown # type: ignore
            # テキストを取得
            with open(filename, "r" ,encoding=encoding, errors='ignore') as f:
                text_data = f.read()
                md = markdown(text_data)
                soup = BeautifulSoup(md, "html.parser")
            result = soup.get_text()
        else:
            # その他のtext/*の場合
            with open(filename, "r", encoding=encoding, errors='ignore') as f:
                result = f.read()
            
        return result

import sys
import datetime

import openpyxl # type: ignore

class ExcelUtil:

    @staticmethod
    def export_to_excel(filePath, data):
        # Workbookオブジェクトを生成
        wb = openpyxl.Workbook()
        # アクティブなシートを取得
        ws = wb.active
        # シート名を設定
        ws.title = "Sheet1"
        # データを書き込む
        for row in data:
            ws.append(row)
        # ファイルを保存
        wb.save(filePath)
        
    @staticmethod
    def import_from_excel(filePath):
        # Workbookオブジェクトを生成
        wb = openpyxl.load_workbook(filePath)
        # アクティブなシートを取得
        ws = wb.active
        # データを取得
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append(row)
        
        return data

    # application/vnd.openxmlformats-officedocument.spreadsheetml.sheetのファイルを読み込んで文字列として返す関数
    @staticmethod
    def extract_text_from_sheet(filename:str, sheet_name:str=""):
        import openpyxl
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        wb = openpyxl.load_workbook(filename)
        for sheet in wb:
            # シート名が指定されている場合はそのシートのみ処理
            if sheet_name and sheet.title != sheet_name:
                continue
            for row in sheet.iter_rows(values_only=True):
                # 1行分のデータを格納するリスト
                cells = []
                for cell in row:
                    # cell.valueがNoneの場合はcontinue
                    if cell is None:
                        continue
                    # cell.valueがdatetime.datetimeの場合はisoformat()で文字列に変換
                    if isinstance(cell, datetime.datetime):
                        cells.append(cell.isoformat())
                    else:
                        cells.append(str(cell))
                    
                output.write("\t".join(cells))
                output.write("\n")
        
        return output.getvalue()

    # excelのシート名一覧を取得する関数
    @staticmethod
    def get_sheet_names(filename):
        import openpyxl
        wb = openpyxl.load_workbook(filename)
        return wb.sheetnames

