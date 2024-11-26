import json

from typing import Annotated, Callable

# wikipedia
import wikipedia
# selenium
from selenium import webdriver
from selenium.webdriver.edge.service import Service
from selenium.webdriver.edge.options import Options
from webdriver_manager.microsoft import EdgeChromiumDriverManager

class AutoGenToolGenerator:

    @classmethod        
    def create_default_tools(cls) -> dict[str, tuple[Callable, str]]:
        # デフォルトのツールを生成
        # 関数名と関数生成関数のペアを保持する辞書
        tools: dict[str, tuple[Callable, str]] = {
            "search_wikipedia_ja": cls.__create_search_wikipedia_ja(),
            "list_files_in_directory": cls.__create_list_files_in_directory(),
            "extract_text_from_file": cls.__create_extract_text_from_file(),
            "extract_webpage": cls.__create_extract_webpage_function(),
            "check_file": cls.__create_check_file_function(),
            "save_text_file": cls.__create_save_text_file_function(),
            "search_duckduckgo": cls.__create_search_duckduckgo(),
            "get_current_time": cls.__create_get_current_time()
        }
        return tools

    @classmethod
    def __create_search_wikipedia_ja(cls) -> tuple[Callable[[str, int], list[str]], str]:
        def search_wikipedia_ja(query: Annotated[str, "String to search for"], lang: Annotated[str, "Language of Wikipedia"], num_results: Annotated[int, "Maximum number of results to display"]) -> list[str]:

            # Use the Japanese version of Wikipedia
            wikipedia.set_lang(lang)
            
            # Retrieve search results
            search_results = wikipedia.search(query, results=num_results)
            
            result_texts = []
            # Display the top results
            for i, title in enumerate(search_results):
            
                print(f"Result {i + 1}: {title}")
                try:
                    # Retrieve the content of the page
                    page = wikipedia.page(title)
                    print(page.content[:500])  # Display the first 500 characters
                    text = f"Title:\n{title}\n\nContent:\n{page.content}\n"
                    result_texts.append(text)
                except wikipedia.exceptions.DisambiguationError as e:
                    print(f"Disambiguation: {e.options}")
                except wikipedia.exceptions.PageError:
                    print("Page not found.")
                print("\n" + "-"*50 + "\n")
            return result_texts
        return search_wikipedia_ja, "This function searches Wikipedia with the specified keywords and returns related articles."

    @classmethod
    def __create_list_files_in_directory(cls) -> tuple[Callable[[str], list[str]], str]:
        def list_files_in_directory(directory_path: Annotated[str, "Directory path"]) -> list[str]:
            import os
            files = os.listdir(directory_path)
            return files
        return list_files_in_directory, "This function returns a list of files in the specified directory."

    @classmethod
    def __create_extract_text_from_file(cls) -> tuple[Callable[[str], str]]:
        # Generates a function to extract text from a file
        def extract_file(file_path: Annotated[str, "File path"]) -> str:
            # Extract text from a temporary file
            text = FileUtil.extract_text_from_file(file_path)
            return text
        return extract_file, "This function extracts text from the specified file."

    @classmethod
    def __create_check_file_function(cls) -> tuple[Callable[[str], bool], str]:
        def check_file(file_path: Annotated[str, "File path"]) -> bool:
            # Check if the file exists
            import os
            check_file = os.path.exists(file_path)
            return check_file
        
        return check_file, "This function checks if the specified file exists."

    @classmethod
    def __create_extract_webpage_function(cls) -> tuple[Callable[[str], list[str]], str]:
        def extract_webpage(url: Annotated[str, "URL of the web page to extract text and links from"]) -> Annotated[tuple[str, list[tuple[str, str]]], "Page text, list of links (href attribute and link text from <a> tags)"]:
            driver = cls.__create_web_driver()
            # Wait for the page to fully load (set explicit wait conditions if needed)
            driver.implicitly_wait(10)
            # Retrieve HTML of the web page and extract text and links
            driver.get(url)
            # Get the entire HTML of the page
            page_html = driver.page_source

            from bs4 import BeautifulSoup
            soup = BeautifulSoup(page_html, "html.parser")
            text = soup.get_text()
            # Retrieve href attribute and text from <a> tags
            urls: list[tuple[str, str]] = [(a.get("href"), a.get_text()) for a in soup.find_all("a")]
            driver.close()
            return text, urls
        return extract_webpage, "This function extracts text and links from the specified URL of a web page."

    @classmethod
    def __create_web_driver(cls):
        # ヘッドレスモードのオプションを設定
        edge_options = Options()
        edge_options.add_argument("--headless")
        edge_options.add_argument("--disable-gpu")
        edge_options.add_argument("--no-sandbox")
        edge_options.add_argument("--disable-dev-shm-usage")

        # Edgeドライバをセットアップ
        driver = webdriver.Edge(service=Service(EdgeChromiumDriverManager().install()), options=edge_options)
        return driver

    # Generate a function to search with the specified keywords using the DuckDuckGo API and return related articles
    @classmethod
    def __create_search_duckduckgo(cls) -> tuple[Callable, str]:
        from duckduckgo_search import DDGS
        ddgs = DDGS()
        def search_duckduckgo(query: Annotated[str, "String to search for"], site: Annotated[str, "Site to search within. Leave blank if no site is specified"], num_results: Annotated[int, "Maximum number of results to display"]) -> Annotated[list[tuple[str, str, str]], "(Title, URL, Body) list"]:
            try:
                # Retrieve DuckDuckGo search results
                if site:
                    query = f"{query} site:{site}"

                print(f"Search query: {query}")

                results_dict = ddgs.text(
                    keywords=query,            # Search keywords
                    region='jp-jp',            # Region. For Japan: "jp-jp", No specific region: "wt-wt"
                    safesearch='off',          # Safe search OFF->"off", ON->"on", Moderate->"moderate"
                    timelimit=None,            # Time limit. None for no limit, past day->"d", past week->"w", past month->"m", past year->"y"
                    max_results=num_results    # Number of results to retrieve
                )

                results = []
                for result in results_dict:
                    # title, href, body
                    title = result.get("title", "")
                    href = result.get("href", "")
                    body = result.get("body", "")
                    print(f'Title: {title}, URL: {href}, Body: {body[:100]}')
                    results.append((title, href, body))

                return results
            except Exception as e:
                print(e)
                import traceback
                traceback.print_exc()
                return []
            
        return search_duckduckgo, "This function searches DuckDuckGo with the specified keywords and returns related articles."

    @classmethod
    def __create_save_text_file_function(cls) -> tuple[Callable[[str, str, str], None], str]:
        def save_text_file(name: Annotated[str, "File name"], dirname: Annotated[str, "Directory name"], text: Annotated[str, "Text data to save"]) -> Annotated[bool, "Save result"]:
            # Save in the specified directory
            try:
                import os
                if not os.path.exists(dirname):
                    os.makedirs(dirname)
                path = os.path.join(dirname, name)
                with open(path, "w", encoding="utf-8") as f:
                    f.write(text)
                # Check if the file exists
                return os.path.exists(path)
            except Exception as e:
                print(e)
                return False

        return save_text_file, "This function saves text data as a file."

    # Generate a function that returns the current time in the format yyyy/mm/dd (Day) hh:mm:ss
    @classmethod
    def __create_get_current_time(cls) -> tuple[Callable, str]:
        from datetime import datetime
        def get_current_time() -> str:
            now = datetime.now()
            return now.strftime("%Y/%m/%d (%a) %H:%M:%S")
        return get_current_time, "This function returns the current time in the format yyyy/mm/dd (Day) hh:mm:ss."

from magika import Magika  # type: ignore
import chardet

class FileUtil:

    @classmethod
    def identify_type(cls, filename):
        m = Magika()
        # ファイルのbyte列を取得
        # アクセスできない場合は例外をキャッチ
        try:
            with open(filename, "rb") as f:
                # 1KB読み込む
                byte_data = f.read(8192)
        except Exception as e:
            print(e)
            return None, None

        # ファイルの種類を判定
        res = m.identify_bytes(byte_data)
        # エンコーディング判定
        encoding_dic = chardet.detect(byte_data)
        encoding = encoding_dic["encoding"]
        if encoding is None or encoding == "UTF-8":
            encoding = "utf-8"
        elif encoding == "SHIFT_JIS":
            encoding = "cp932"
            
        return res, encoding

    @classmethod
    def get_mime_type(cls, filename):
        res, encoding = cls.identify_type(filename)
        if res is None:
            return None
        return res.output.mime_type

    @classmethod
    def extract_text_from_file(cls, filename):
        res, encoding = cls.identify_type(filename)
        if res is None:
            return None
        print(res.output.mime_type)
        
        if res.output.mime_type.startswith("text/"):
            return cls.process_text(filename, res, encoding)

        # application/pdf
        elif res.output.mime_type == "application/pdf":
            return cls.process_pdf(filename)
            
        # application/vnd.openxmlformats-officedocument.spreadsheetml.sheet
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return ExcelUtil.extract_text_from_sheet(filename)
            
        # application/vnd.openxmlformats-officedocument.wordprocessingml.document
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return cls.process_docx(filename)
            
        # application/vnd.openxmlformats-officedocument.presentationml.presentation
        elif res.output.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return cls.process_pptx(filename)
        else:
            print("Unsupported file type: " + res.output.mime_type)



    # application/pdfのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_pdf(cls, filename):
        from pdfminer.high_level import extract_text
        text = extract_text(filename)
        return text


    # application/vnd.openxmlformats-officedocument.wordprocessingml.documentのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_docx(cls, filename):
        import docx
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        doc = docx.Document(filename)
        for para in doc.paragraphs:
            output.write(para.text)
            output.write("\n")
            
        return output.getvalue()

    # application/vnd.openxmlformats-officedocument.presentationml.presentationのファイルを読み込んで文字列として返す関数
    @classmethod
    def process_pptx(cls, filename):
        import pptx
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        prs = pptx.Presentation(filename)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output.write(shape.text)
                    output.write("\n")
        
        return output.getvalue()

    # text/*のファイルを読み込んで文字列として返す関数
    @classmethod
    def process_text(cls, filename, res, encoding):
        result = ""
        if res.output.mime_type == "text/html":
            # text/htmlの場合
            from bs4 import BeautifulSoup
            # テキストを取得
            with open(filename, "rb") as f:
                text_data = f.read()
                soup = BeautifulSoup(text_data, "html.parser")
            result = soup.get_text()

        elif res.output.mime_type == "text/xml":
            # text/xmlの場合
            from bs4 import BeautifulSoup
            # テキストを取得
            with open(filename, "rb") as f:
                text_data = f.read()
                soup = BeautifulSoup(text_data, features="xml")
            result = soup.get_text()

        elif res.output.mime_type == "text/markdown":
            # markdownの場合
            from bs4 import BeautifulSoup
            from markdown import markdown
            # テキストを取得
            with open(filename, "r" ,encoding=encoding) as f:
                text_data = f.read()
                md = markdown(text_data)
                soup = BeautifulSoup(md, "html.parser")
            result = soup.get_text()
        else:
            # その他のtext/*の場合
            with open(filename, "r", encoding=encoding) as f:
                result = f.read()
            
        return result

import openpyxl
import datetime

class ExcelUtil:

    @staticmethod
    def export_to_excel(filePath, data):
        # Workbookオブジェクトを生成
        wb = openpyxl.Workbook()
        # アクティブなシートを取得
        ws = wb.active
        # シート名を設定
        ws.title = "Sheet1"
        # データを書き込む
        for row in data:
            ws.append(row)
        # ファイルを保存
        wb.save(filePath)
        
    @staticmethod
    def import_from_excel(filePath):
        # Workbookオブジェクトを生成
        wb = openpyxl.load_workbook(filePath)
        # アクティブなシートを取得
        ws = wb.active
        # データを取得
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append(row)
        
        return data

    # application/vnd.openxmlformats-officedocument.spreadsheetml.sheetのファイルを読み込んで文字列として返す関数
    @staticmethod
    def extract_text_from_sheet(filename:str, sheet_name:str=""):
        import openpyxl
        from io import StringIO
        # 出力用のストリームを作成
        output = StringIO()
        wb = openpyxl.load_workbook(filename)
        for sheet in wb:
            # シート名が指定されている場合はそのシートのみ処理
            if sheet_name and sheet.title != sheet_name:
                continue
            for row in sheet.iter_rows(values_only=True):
                # 1行分のデータを格納するリスト
                cells = []
                for cell in row:
                    # cell.valueがNoneの場合はcontinue
                    if cell is None:
                        continue
                    # cell.valueがdatetime.datetimeの場合はisoformat()で文字列に変換
                    if isinstance(cell, datetime.datetime):
                        cells.append(cell.isoformat())
                    else:
                        cells.append(str(cell))
                    
                output.write("\t".join(cells))
                output.write("\n")
        
        return output.getvalue()

    # excelのシート名一覧を取得する関数
    @staticmethod
    def get_sheet_names(filename):
        import openpyxl
        wb = openpyxl.load_workbook(filename)
        return wb.sheetnames
